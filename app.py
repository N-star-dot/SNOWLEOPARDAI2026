import streamlit as st
import pandas as pd
import sqlite3
import json
import base64
import io
import html
import agent
from datetime import datetime

# --- 1. UI SETUP ---
st.set_page_config(page_title="Acuman AI", page_icon="🌊", layout="wide")

# --- 2. MEMORY & STATE ---
if "messages" not in st.session_state:
    st.session_state.messages = []

if "active_mode" not in st.session_state:
    st.session_state.active_mode = "Cloud (Snow Leopard - IMDb)"

# --- NEW: Initialize the dynamic UI variables ---
if "confirmed_mode" not in st.session_state:
    st.session_state.confirmed_mode = st.session_state.active_mode
if "temperature" not in st.session_state:
    st.session_state.temperature = 0.2
if "contact_target" not in st.session_state:
    st.session_state.contact_target = ""
if "voice_transcript" not in st.session_state:
    st.session_state.voice_transcript = None
if "last_ingested_file_id" not in st.session_state:
    st.session_state.last_ingested_file_id = None
if "uploaded_image_b64" not in st.session_state:
    st.session_state.uploaded_image_b64 = None
if "uploaded_video_bytes" not in st.session_state:
    st.session_state.uploaded_video_bytes = None
if "run_proactive_eda" not in st.session_state:
    st.session_state.run_proactive_eda = False

# THE AMNESIA FIX: Check the hard drive for the database on startup!
if "dataset_columns" not in st.session_state:
    try:
        temp_conn = sqlite3.connect('hackathon_database.db')
        cols_df = pd.read_sql_query("PRAGMA table_info(user_data)", temp_conn)
        if not cols_df.empty:
            st.session_state.dataset_columns = ", ".join(cols_df['name'].tolist())
        else:
            st.session_state.dataset_columns = "No local data uploaded yet."
    except Exception:
        st.session_state.dataset_columns = "No local data uploaded yet."
# --- DYNAMIC THEME ENGINE ---
is_cloud = "Cloud" in st.session_state.confirmed_mode

# Shared font import
FONT_IMPORT = "@import url('https://fonts.googleapis.com/css2?family=Shippori+Mincho:wght@400;500;600;700&family=Zen+Kaku+Gothic+New:wght@300;400;500;700&display=swap');"

if is_cloud:
    # ❄️ SNOW LEOPARD — glacial palette, ma spacing, frost glass, breathe pulse
    ACCENT_U = "#8fa8c0"
    theme_css = f"""
    <style>
    {FONT_IMPORT}

    /* ── Base ── */
    html, body, [data-testid="stAppViewContainer"] {{
        font-family: 'Zen Kaku Gothic New', sans-serif;
        background: #f0f4f8;
        background-attachment: fixed;
        color: #1a2a38;
    }}
    [data-testid="stAppViewContainer"] {{
        background: linear-gradient(160deg, #f0f4f8 0%, #dde8f2 50%, #cfd9e8 100%);
        background-attachment: fixed;
    }}
    /* subtle frost overlay — breathes like fog on glass */
    [data-testid="stAppViewContainer"]::after {{
        content: '';
        position: fixed;
        inset: 0;
        background: radial-gradient(ellipse at 30% 40%, rgba(143,168,192,0.07) 0%, transparent 65%),
                    radial-gradient(ellipse at 75% 70%, rgba(58,80,104,0.05) 0%, transparent 55%);
        pointer-events: none;
        z-index: 0;
        animation: frost-breathe 8s ease-in-out infinite;
    }}
    @keyframes frost-breathe {{
        0%,100% {{ opacity: 0.6; }}
        50%      {{ opacity: 1; }}
    }}
    [data-testid="stHeader"] {{ background: transparent; }}

    /* ── Wordmark ── */
    .acuman-wordmark {{
        font-family: 'Shippori Mincho', serif;
        font-size: 2rem;
        font-weight: 600;
        letter-spacing: 0.04em;
        color: #1a2a38;
        line-height: 1.2;
        margin: 0.5rem 0 0.25rem 0;
    }}
    .acuman-wordmark .u-pivot {{
        color: {ACCENT_U};
        font-style: italic;
    }}
    .acuman-sub {{
        font-family: 'Zen Kaku Gothic New', sans-serif;
        font-size: 0.78rem;
        font-weight: 300;
        letter-spacing: 0.18em;
        text-transform: uppercase;
        color: #6a8aa8;
        margin-top: 0;
    }}

    /* ── Headings ── */
    h1, h2, h3 {{
        font-family: 'Shippori Mincho', serif !important;
        color: #1a2a38 !important;
        font-weight: 600 !important;
        letter-spacing: 0.02em !important;
    }}
    h2 {{ font-size: 1.35rem !important; }}
    h3 {{ font-size: 1.1rem !important; }}

    /* ── Sidebar ── */
    [data-testid="stSidebar"] {{
        background: linear-gradient(180deg, #1a2a38 0%, #0f1e2e 100%) !important;
        border-right: 0.5px solid rgba(143,168,192,0.25) !important;
        padding: 1.5rem 0 !important;
    }}
    [data-testid="stSidebar"] p,
    [data-testid="stSidebar"] label,
    [data-testid="stSidebar"] div[data-testid="stMarkdownContainer"] > p {{
        color: #c8dcea !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}
    [data-testid="stSidebar"] h1,
    [data-testid="stSidebar"] h2,
    [data-testid="stSidebar"] h3 {{
        font-family: 'Shippori Mincho', serif !important;
        color: #ddeeff !important;
        letter-spacing: 0.06em !important;
    }}
    [data-testid="stSidebar"] hr {{
        border: none;
        border-top: 0.5px solid rgba(143,168,192,0.18);
        margin: 1.2rem 0;
    }}

    /* ── Chat Bubbles — frosted glass ── */
    [data-testid="stChatMessage"] {{
        background: rgba(240,244,248,0.72);
        backdrop-filter: blur(18px);
        -webkit-backdrop-filter: blur(18px);
        border-radius: 12px;
        padding: 20px 24px;
        box-shadow: 0 2px 16px rgba(58,80,104,0.08), inset 0 0.5px 0 rgba(255,255,255,0.9);
        border: 0.5px solid rgba(143,168,192,0.28);
        color: #1a2a38;
        margin-bottom: 1rem;
        transition: box-shadow 0.4s ease, border-color 0.4s ease;
    }}
    [data-testid="stChatMessage"]:hover {{
        box-shadow: 0 4px 24px rgba(58,80,104,0.13), inset 0 0.5px 0 rgba(255,255,255,1);
        border-color: rgba(143,168,192,0.45);
    }}

    /* ── Chat Input ── */
    [data-testid="stChatInputContainer"] > div,
    .stChatInput > div {{
        background: rgba(240,244,248,0.8) !important;
        border: 0.5px solid rgba(143,168,192,0.35) !important;
        border-radius: 12px !important;
        backdrop-filter: blur(12px);
    }}
    .stChatInput textarea {{
        color: #1a2a38 !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}

    /* ── Expanders ── */
    [data-testid="stExpander"] {{
        background: rgba(240,244,248,0.6);
        backdrop-filter: blur(10px);
        border: 0.5px solid rgba(143,168,192,0.22);
        border-radius: 12px;
    }}

    /* ── Buttons ── */
    .stButton > button {{
        background: #3a5068 !important;
        color: #ddeeff !important;
        border: 0.5px solid rgba(143,168,192,0.4) !important;
        border-radius: 8px !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
        font-weight: 500 !important;
        letter-spacing: 0.06em !important;
        transition: all 0.35s ease !important;
    }}
    .stButton > button:hover {{
        background: #4a6880 !important;
        box-shadow: 0 0 18px rgba(143,168,192,0.22) !important;
        transform: translateY(-1px) !important;
    }}
    .stDownloadButton > button {{
        background: #3a5068 !important;
        color: #ddeeff !important;
        border-radius: 8px !important;
        border: 0.5px solid rgba(143,168,192,0.3) !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}

    /* ── Slider ── */
    .stSlider > div > div > div {{
        background-color: {ACCENT_U} !important;
    }}

    /* ── Text ── */
    p, label, .stMarkdown, .stCaption, small,
    [data-testid="stMarkdownContainer"] > p,
    [data-testid="stMarkdownContainer"] > span {{
        color: #2a3c50 !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}
    
    /* Ensure Streamlit icons (which are often empty spans or use specific font families) are not overridden */
    span.st-emotion-cache-1vt4ygl, 
    span.stIcon {{
        font-family: inherit !important;
    }}
    [data-testid="stAppViewContainer"] .stMarkdown h1,
    [data-testid="stAppViewContainer"] .stMarkdown h2,
    [data-testid="stAppViewContainer"] .stMarkdown h3 {{
        font-family: 'Shippori Mincho', serif !important;
        color: #1a2a38 !important;
    }}

    /* ── Divider ── */
    hr {{
        border: none;
        border-top: 0.5px solid rgba(143,168,192,0.2);
        margin: 1.8rem 0;
    }}
    
    /* ── Ambient Particles (Snow) ── */
    .particle-container {{
        position: fixed; inset: 0; pointer-events: none; z-index: 0; overflow: hidden;
    }}
    .particle-snow {{
        position: absolute; width: 6px; height: 6px;
        background-image: url("data:image/svg+xml,%3Csvg width='24' height='24' viewBox='0 0 24 24' fill='none' stroke='%238fa8c0' stroke-width='1.5' stroke-linecap='round' stroke-linejoin='round' xmlns='http://www.w3.org/2000/svg'%3E%3Cline x1='12' y1='2' x2='12' y2='22'%3E%3C/line%3E%3Cline x1='17' y1='5' x2='7' y2='19'%3E%3C/line%3E%3Cline x1='5' y1='17' x2='19' y2='7'%3E%3C/line%3E%3C/svg%3E");
        background-size: contain; opacity: 0.25;
        animation: snow-drift linear infinite;
    }}
    @keyframes snow-drift {{
        0% {{ transform: translateY(-5vh) translateX(0) rotate(0deg); opacity: 0; }}
        10% {{ opacity: 0.4; }}
        90% {{ opacity: 0.4; }}
        100% {{ transform: translateY(105vh) translateX(20px) rotate(180deg); opacity: 0; }}
    }}
    
    /* ── Radio Sigil (Snowflake) ── */
    div[role="radiogroup"] label[data-baseweb="radio"] div:first-child {{
        background-color: transparent !important;
        border: 1px solid rgba(143,168,192,0.4) !important;
    }}
    div[role="radiogroup"] label[data-baseweb="radio"] input:checked + div {{
        background-color: {ACCENT_U} !important;
        border-color: {ACCENT_U} !important;
        background-image: url("data:image/svg+xml,%3Csvg width='12' height='12' viewBox='0 0 24 24' fill='none' stroke='white' stroke-width='2.5' stroke-linecap='round' stroke-linejoin='round' xmlns='http://www.w3.org/2000/svg'%3E%3Cline x1='12' y1='3' x2='12' y2='21'%3E%3C/line%3E%3Cline x1='19' y1='7.5' x2='5' y2='16.5'%3E%3C/line%3E%3Cline x1='5' y1='7.5' x2='19' y2='16.5'%3E%3C/line%3E%3C/svg%3E") !important;
        background-position: center !important;
        background-repeat: no-repeat !important;
    }}
    div[role="radiogroup"] label[data-baseweb="radio"] input:checked + div div {{ display: none; }}
    
    /* ── Slider Personality (Ice) ── */
    div[data-testid="stThumbValue"] {{
        background: {ACCENT_U} !important;
        color: #fff !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}
    div[role="slider"] {{
        background-color: #f0f4f8 !important;
        border: 2px solid {ACCENT_U} !important;
        box-shadow: 0 0 8px rgba(143,168,192,0.4) !important;
        background-image: url("data:image/svg+xml,%3Csvg width='12' height='12' viewBox='0 0 24 24' fill='none' stroke='%238fa8c0' stroke-width='2' stroke-linecap='round' stroke-linejoin='round' xmlns='http://www.w3.org/2000/svg'%3E%3Cline x1='12' y1='4' x2='12' y2='20'%3E%3C/line%3E%3Cline x1='18' y1='8' x2='6' y2='16'%3E%3C/line%3E%3Cline x1='6' y1='8' x2='18' y2='16'%3E%3C/line%3E%3C/svg%3E") !important;
        background-position: center !important;
        background-repeat: no-repeat !important;
        transition: transform 0.2s ease !important;
    }}
    div[role="slider"]:hover {{ transform: scale(1.15) !important; }}
    
    /* ── Mode Transition Overlay ── */
    .mode-transition-overlay {{
        position: fixed; inset: 0; background: #dde8f2; z-index: 99999; pointer-events: none;
        animation: fade-out 600ms ease-out forwards;
    }}
    @keyframes fade-out {{ from {{ opacity: 1; }} to {{ opacity: 0; visibility: hidden; }} }}
    
    /* ── Input Bar Frost ── */
    [data-testid="stChatInputContainer"] {{
        background: linear-gradient(0deg, rgba(240,244,248,0.95) 0%, rgba(240,244,248,0.7) 100%) !important;
        backdrop-filter: blur(12px);
        -webkit-backdrop-filter: blur(12px);
        border-top: 0.5px solid rgba(143,168,192,0.2) !important;
        padding-top: 1rem !important;
    }}
    [data-testid="stChatInputContainer"] button {{
        color: {ACCENT_U} !important;
    }}
    
    /* ── Enso Drawing Animation ── */
    .enso-draw {{
        stroke-dasharray: 100;
        stroke-dashoffset: 100;
        animation: draw-enso 1.5s ease-in-out infinite alternate;
        transform-origin: center;
        transform: rotate(-90deg);
    }}
    @keyframes draw-enso {{
        0% {{ stroke-dashoffset: 100; }}
        100% {{ stroke-dashoffset: 0; }}
    }}
    </style>
    
    <div class="mode-transition-overlay"></div>
    <div class="particle-container">
        <div class="particle-snow" style="left:15%; animation-duration: 22s; animation-delay: 0s; width:8px; height:8px;"></div>
        <div class="particle-snow" style="left:35%; animation-duration: 30s; animation-delay: -5s; width:5px; height:5px; opacity: 0.15;"></div>
        <div class="particle-snow" style="left:65%; animation-duration: 25s; animation-delay: -12s; width:7px; height:7px;"></div>
        <div class="particle-snow" style="left:85%; animation-duration: 28s; animation-delay: -2s; width:6px; height:6px; opacity: 0.2;"></div>
        <div class="particle-snow" style="left:45%; animation-duration: 35s; animation-delay: -18s; width:4px; height:4px; opacity: 0.1;"></div>
        <div class="particle-snow" style="left:5%; animation-duration: 20s; animation-delay: -8s; width:9px; height:9px; opacity: 0.3;"></div>
        <div class="particle-snow" style="left:55%; animation-duration: 27s; animation-delay: -15s; width:7px; height:7px;"></div>
        <div class="particle-snow" style="left:92%; animation-duration: 24s; animation-delay: -4s; width:6px; height:6px; opacity: 0.25;"></div>
    </div>
    """
    theme_name = " "
    theme_sub = " "
else:
    # 🍵 LOCAL — washi paper, hinoki warmth, candle flicker, paper grain
    ACCENT_U_LOCAL = "#c8a87a"
    theme_css = f"""
    <style>
    {FONT_IMPORT}

    /* SVG paper grain as data URI background */
    html, body, [data-testid="stAppViewContainer"] {{
        font-family: 'Zen Kaku Gothic New', sans-serif;
        color: #2a1d0e;
    }}
    [data-testid="stAppViewContainer"] {{
        background-color: #f5f0e8;
        background-image: url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' width='400' height='400'%3E%3Cfilter id='grain'%3E%3CfeTurbulence type='fractalNoise' baseFrequency='0.75' numOctaves='4' stitchTiles='stitch'/%3E%3CfeColorMatrix type='saturate' values='0'/%3E%3C/filter%3E%3Crect width='400' height='400' filter='url(%23grain)' opacity='0.035'/%3E%3C/svg%3E"),
                          linear-gradient(160deg, #f5f0e8 0%, #ede0cc 45%, #e8d9c0 100%);
        background-attachment: fixed;
    }}
    /* candle flicker on the entire page atmosphere */
    [data-testid="stAppViewContainer"]::after {{
        content: '';
        position: fixed;
        inset: 0;
        background: radial-gradient(ellipse at 50% 30%, rgba(200,168,122,0.06) 0%, transparent 70%);
        pointer-events: none;
        z-index: 0;
        animation: candle-flicker 3s ease-in-out infinite;
    }}
    @keyframes candle-flicker {{
        0%,100% {{ opacity: 0.85; }}
        20%      {{ opacity: 1; }}
        40%      {{ opacity: 0.92; }}
        60%      {{ opacity: 0.88; }}
        80%      {{ opacity: 0.97; }}
    }}
    [data-testid="stHeader"] {{ background: transparent; }}

    /* ── Wordmark ── */
    .acuman-wordmark {{
        font-family: 'Shippori Mincho', serif;
        font-size: 2rem;
        font-weight: 600;
        letter-spacing: 0.04em;
        color: #2a1d0e;
        line-height: 1.2;
        margin: 0.5rem 0 0.25rem 0;
    }}
    .acuman-wordmark .u-pivot {{
        color: {ACCENT_U_LOCAL};
        font-style: italic;
    }}
    .acuman-sub {{
        font-family: 'Zen Kaku Gothic New', sans-serif;
        font-size: 0.78rem;
        font-weight: 300;
        letter-spacing: 0.18em;
        text-transform: uppercase;
        color: #8a6a48;
        margin-top: 0;
    }}

    /* ── Headings ── */
    h1, h2, h3 {{
        font-family: 'Shippori Mincho', serif !important;
        color: #2a1d0e !important;
        font-weight: 600 !important;
        letter-spacing: 0.02em !important;
    }}
    h2 {{ font-size: 1.35rem !important; }}
    h3 {{ font-size: 1.1rem !important; }}

    /* ── Sidebar ── */
    [data-testid="stSidebar"] {{
        background: linear-gradient(180deg, #f0e8d8 0%, #e8dcc8 100%) !important;
        border-right: 0.5px solid rgba(200,168,122,0.28) !important;
        padding: 1.5rem 0 !important;
    }}
    [data-testid="stSidebar"] p,
    [data-testid="stSidebar"] label,
    [data-testid="stSidebar"] div[data-testid="stMarkdownContainer"] > p {{
        color: #3a2a18 !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}
    [data-testid="stSidebar"] h1,
    [data-testid="stSidebar"] h2,
    [data-testid="stSidebar"] h3 {{
        font-family: 'Shippori Mincho', serif !important;
        color: #2a1d0e !important;
        letter-spacing: 0.06em !important;
    }}
    [data-testid="stSidebar"] hr {{
        border: none;
        border-top: 0.5px solid rgba(200,168,122,0.22);
        margin: 1.2rem 0;
    }}

    /* ── Chat Bubbles — washi warmth ── */
    [data-testid="stChatMessage"] {{
        background: rgba(245,240,232,0.88);
        backdrop-filter: blur(10px);
        -webkit-backdrop-filter: blur(10px);
        border-radius: 18px;
        padding: 20px 24px;
        box-shadow: 0 4px 24px rgba(90,50,10,0.06), inset 0 0.5px 0 rgba(255,255,255,0.85);
        border: 0.5px solid rgba(200,168,122,0.28);
        color: #2a1d0e;
        margin-bottom: 1rem;
        transition: box-shadow 0.4s ease;
    }}
    [data-testid="stChatMessage"]:hover {{
        box-shadow: 0 6px 28px rgba(90,50,10,0.10), inset 0 0.5px 0 rgba(255,255,255,1);
        border-color: rgba(200,168,122,0.42);
    }}

    /* ── Chat Input ── */
    [data-testid="stChatInputContainer"] > div,
    .stChatInput > div {{
        background: rgba(245,240,232,0.9) !important;
        border: 0.5px solid rgba(200,168,122,0.35) !important;
        border-radius: 14px !important;
    }}
    .stChatInput textarea {{
        color: #2a1d0e !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}

    /* ── Expanders ── */
    [data-testid="stExpander"] {{
        background: rgba(245,240,232,0.65);
        border: 0.5px solid rgba(200,168,122,0.22);
        border-radius: 16px;
    }}

    /* ── Buttons ── */
    .stButton > button {{
        background: #5c3d1e !important;
        color: #f5f0e8 !important;
        border: 0.5px solid rgba(200,168,122,0.35) !important;
        border-radius: 10px !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
        font-weight: 500 !important;
        letter-spacing: 0.05em !important;
        transition: all 0.35s ease !important;
    }}
    .stButton > button:hover {{
        background: #6e4a28 !important;
        box-shadow: 0 4px 16px rgba(92,61,30,0.2) !important;
        transform: translateY(-1px) !important;
    }}
    .stDownloadButton > button {{
        background: #5c3d1e !important;
        color: #f5f0e8 !important;
        border-radius: 10px !important;
        border: 0.5px solid rgba(200,168,122,0.3) !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}

    /* ── Slider ── */
    .stSlider > div > div > div {{
        background-color: {ACCENT_U_LOCAL} !important;
    }}

    /* ── Text ── */
    p, label, .stMarkdown, .stCaption, small,
    [data-testid="stMarkdownContainer"] > p,
    [data-testid="stMarkdownContainer"] > span {{
        color: #3a2a18 !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}
    [data-testid="stAppViewContainer"] .stMarkdown h1,
    [data-testid="stAppViewContainer"] .stMarkdown h2,
    [data-testid="stAppViewContainer"] .stMarkdown h3 {{
        font-family: 'Shippori Mincho', serif !important;
        color: #2a1d0e !important;
    }}

    /* ── Divider ── */
    hr {{
        border: none;
        border-top: 0.5px solid rgba(200,168,122,0.22);
        margin: 1.8rem 0;
    }}
    
    /* ── Ambient Particles (Smoke Wisps) ── */
    .particle-container {{
        position: fixed; inset: 0; pointer-events: none; z-index: 0; overflow: hidden;
    }}
    .smoke-wisp {{
        position: absolute; bottom: -10vh; width: 2px; height: 30vh;
        background: linear-gradient(0deg, transparent 0%, rgba(200,168,122,0.08) 50%, transparent 100%);
        filter: blur(8px);
        animation: smoke-rise linear infinite;
    }}
    @keyframes smoke-rise {{
        0% {{ transform: translateY(0) rotate(0deg) scaleX(1); opacity: 0; }}
        20% {{ opacity: 1; }}
        80% {{ opacity: 1; }}
        100% {{ transform: translateY(-110vh) rotate(10deg) scaleX(3); opacity: 0; }}
    }}
    
    /* ── Radio Sigil (Leaf) ── */
    div[role="radiogroup"] label[data-baseweb="radio"] div:first-child {{
        background-color: transparent !important;
        border: 1px solid rgba(200,168,122,0.4) !important;
    }}
    div[role="radiogroup"] label[data-baseweb="radio"] input:checked + div {{
        background-color: {ACCENT_U_LOCAL} !important;
        border-color: {ACCENT_U_LOCAL} !important;
        background-image: url("data:image/svg+xml,%3Csvg width='12' height='12' viewBox='0 0 24 24' fill='none' stroke='white' stroke-width='2.5' stroke-linecap='round' stroke-linejoin='round' xmlns='http://www.w3.org/2000/svg'%3E%3Cpath d='M11 20A7 7 0 0 1 9.8 6.1C15.5 5 17 4.48 19 2c1 2 2 4.18 2 8 0 5.5-4.78 10-10 10Z'%3E%3C/path%3E%3Cpath d='M2 22l10-10'%3E%3C/path%3E%3C/svg%3E") !important;
        background-position: center !important;
        background-repeat: no-repeat !important;
    }}
    div[role="radiogroup"] label[data-baseweb="radio"] input:checked + div div {{ display: none; }}
    
    /* ── Slider Personality (Amber Drop) ── */
    div[data-testid="stThumbValue"] {{
        background: {ACCENT_U_LOCAL} !important;
        color: #fff !important;
        font-family: 'Zen Kaku Gothic New', sans-serif !important;
    }}
    div[role="slider"] {{
        background-color: #f5f0e8 !important;
        border: 2px solid {ACCENT_U_LOCAL} !important;
        box-shadow: 0 0 8px rgba(200,168,122,0.3) !important;
        border-top-left-radius: 50% !important;
        border-top-right-radius: 50% !important;
        border-bottom-left-radius: 50% !important;
        border-bottom-right-radius: 4px !important;
        transform: rotate(-45deg) !important;
        transition: all 0.2s ease !important;
    }}
    div[role="slider"]:hover {{ transform: rotate(-45deg) scale(1.15) !important; box-shadow: 2px 2px 10px rgba(200,168,122,0.4) !important; }}
    div[role="slider"]:focus {{ transform: rotate(-45deg) scale(1.15) !important; }}
    
    /* ── Mode Transition Overlay ── */
    .mode-transition-overlay {{
        position: fixed; inset: 0; background: #ede0cc; z-index: 99999; pointer-events: none;
        animation: fade-out 600ms ease-out forwards;
    }}
    @keyframes fade-out {{ from {{ opacity: 1; }} to {{ opacity: 0; visibility: hidden; }} }}
    
    /* ── Input Bar Warmth ── */
    [data-testid="stChatInputContainer"] {{
        background: linear-gradient(0deg, rgba(245,240,232,1) 0%, rgba(245,240,232,0.85) 100%) !important;
        backdrop-filter: blur(8px);
        -webkit-backdrop-filter: blur(8px);
        border-top: 0.5px solid rgba(200,168,122,0.25) !important;
        padding-top: 1rem !important;
    }}
    [data-testid="stChatInputContainer"] button {{
        color: {ACCENT_U_LOCAL} !important;
    }}
    
    /* ── Enso Drawing Animation ── */
    .enso-draw {{
        stroke-dasharray: 100;
        stroke-dashoffset: 100;
        animation: draw-enso 1.5s cubic-bezier(0.4, 0, 0.2, 1) infinite alternate;
        transform-origin: center;
        transform: rotate(-90deg);
    }}
    @keyframes draw-enso {{
        0% {{ stroke-dashoffset: 100; }}
        100% {{ stroke-dashoffset: 0; }}
    }}
    </style>
    
    <div class="mode-transition-overlay"></div>
    <div class="particle-container">
        <div class="smoke-wisp" style="left:20%; animation-duration: 45s; animation-delay: 0s;"></div>
        <div class="smoke-wisp" style="left:50%; animation-duration: 55s; animation-delay: -15s; width: 3px; filter: blur(12px); opacity: 0.7;"></div>
        <div class="smoke-wisp" style="left:80%; animation-duration: 65s; animation-delay: -30s;"></div>
    </div>
    """
    theme_name = " "
    theme_sub = " "

st.markdown(theme_css, unsafe_allow_html=True)

# ── Acuman Wordmark ──
if is_cloud:
    u_color = "#8fa8c0"
else:
    u_color = "#c8a87a"

st.markdown(
    f"""
    <div class="acuman-wordmark">Ac<span class="u-pivot">u</span>man</div>
    <p class="acuman-sub">{theme_name} &nbsp;·&nbsp; {theme_sub}</p>
    """,
    unsafe_allow_html=True
)

# ── HAIKU TIME LOGIC ──
current_hour = datetime.now().hour
if 5 <= current_hour < 12:
    haiku = "The mind wakes before the sun"
elif 12 <= current_hour < 17:
    haiku = "Each question, a stone skipped on water"
elif 17 <= current_hour < 21:
    haiku = "Slow thoughts carry the furthest"
else:
    haiku = "Wisdom rests, but does not sleep"

st.markdown(f"<p style='text-align:center; font-family:\"Zen Kaku Gothic New\", sans-serif; font-style:italic; font-weight:300; opacity:0.65; margin-top:-10px; margin-bottom: 2rem; font-size: 0.9rem;'>{haiku}</p>", unsafe_allow_html=True)

# --- 3. MINIMALIST SIDEBAR (Tabs Setup) ---
with st.sidebar:
    st.write("") # Spacer
    
    # Custom CSS for sidebar typography hierarchy
    st.markdown("""
        <style>
        [data-testid="stSidebar"] h3 {
            font-family: 'Shippori Mincho', serif !important;
            font-size: 14px !important;
            letter-spacing: 0.05em !important;
            margin-bottom: -5px !important;
            opacity: 0.9;
        }
        [data-testid="stSidebar"] p, [data-testid="stSidebar"] label {
            font-weight: 300 !important;
            font-size: 13px !important;
        }
        /* Style the tabs */
        .stTabs [data-baseweb="tab-list"] {
            gap: 12px;
        }
        .stTabs [data-baseweb="tab"] {
            padding-left: 8px; padding-right: 8px;
            font-family: 'Zen Kaku Gothic New', sans-serif !important;
            letter-spacing: 0.05em;
        }
        </style>
    """, unsafe_allow_html=True)
    
    # The Tabs to hide complexity — no emojis, clean text
    tab_core, tab_data, tab_tools = st.tabs(["Core", "Data", "Hooks"])

    # TAB 1: NEURAL ROUTING + CREATIVITY SLIDER
    with tab_core:
        st.markdown("### Neural Focus")
        
        selected_mode = st.radio(
            "Select Agent Environment:",
            ["Cloud (Snow Leopard)", "Local (Uploaded CSV - SQLite)"],
            index=0 if "Cloud" in st.session_state.active_mode else 1,
            label_visibility="collapsed"
        )
        st.session_state.active_mode = selected_mode
        
        # Show mode change indicator
        mode_changed = st.session_state.active_mode != st.session_state.confirmed_mode
        
        if mode_changed:
            new_is_cloud = "Cloud" in st.session_state.active_mode
            if new_is_cloud:
                st.info("❄️ Switch to **Snow Leopard** theme?")
            else:
                st.info("🍵 Switch to **Cozy Local** theme?")
            
            if st.button("✅ Confirm Mode Switch", use_container_width=True, key="confirm_mode"):
                st.session_state.confirmed_mode = st.session_state.active_mode
                st.rerun()
        else:
            if is_cloud:
                st.caption("Snow Leopard active — precision.")
            else:
                st.caption("Local mode active — warmth.")
        
        st.divider()
        
        # CREATIVITY SLIDER
        st.markdown("### Intuition")
        st.session_state.temperature = st.slider(
            "Temperature",
            min_value=0.0,
            max_value=1.0,
            value=st.session_state.temperature,
            step=0.05,
            help="Low = strict, clinical reporting. High = creative brainstorming."
        )
        temp_val = st.session_state.temperature
        if temp_val <= 0.3:
            st.caption("Precise & Clinical")
        elif temp_val <= 0.5:
            st.caption("Balanced")
        elif temp_val <= 0.7:
            st.caption("Exploratory")
        else:
            st.caption("Dreaming")

        st.divider()

        # WHISPER VOICE CONTROL
        st.markdown("### Voice")
        try:
            from audio_recorder_streamlit import audio_recorder
            audio_bytes = audio_recorder(
                text="Click to record",
                recording_color="#e74c3c",
                neutral_color="#6c757d",
                icon_size="2x",
                pause_threshold=2.0
            )
            if audio_bytes:
                # Basic check to avoid sending tiny empty clicks (WAV header only)
                if len(audio_bytes) < 1000:
                    st.warning("Audio too short. Please speak longer.")
                else:
                    with st.spinner("Transcribing with Whisper..."):
                        try:
                            audio_file = io.BytesIO(audio_bytes)
                            audio_file.name = "recording.wav"
                            transcription = agent.client.audio.transcriptions.create(
                                file=audio_file,
                                model="whisper-large-v3-turbo",
                                response_format="text",
                            )
                            st.session_state.voice_transcript = transcription
                            st.success(f'Transcribed: "{transcription}"')
                        except Exception as e:
                            st.error(f"Whisper API error: {e}")
        except ImportError:
            st.caption("Install `audio-recorder-streamlit` to enable voice input.")

    # TAB 2: THE LIBRARIAN (Data Ingestion)
    with tab_data:
        st.markdown("### Data Ingestion")
        data_type = st.selectbox("Format", ("CSV (Spreadsheet)", "JSON", "PDF (Document)", "PPTX (Slides)", "DOCX (Document)"), label_visibility="collapsed")
        
        # Set accepted file types based on selection
        accepted_types = ["csv", "json"]
        if data_type == "PDF (Document)":
            accepted_types = ["pdf"]
        elif data_type == "PPTX (Slides)":
            accepted_types = ["pptx"]
        elif data_type == "DOCX (Document)":
            accepted_types = ["docx"]
        
        uploaded_file = st.file_uploader("Drop dataset here", type=accepted_types)
        
        if uploaded_file is not None:
            conn = None
            # Guard to prevent infinite re-processing and reruns
            file_unique_id = f"{uploaded_file.name}_{uploaded_file.size}"
            if st.session_state.get("last_ingested_file_id") != file_unique_id:
                try:
                    conn = sqlite3.connect('hackathon_database.db', timeout=10)
                    MAX_COLUMNS = 500
                    
                    if data_type == "CSV (Spreadsheet)":
                        chunk_size = 10000
                        first_chunk = True
                        progress_text = st.empty()
                        progress_text.text("Processing chunks...")
                        
                        try:
                            # Try robust reading with multiple encodings
                            try:
                                reader = pd.read_csv(uploaded_file, chunksize=chunk_size, on_bad_lines='skip')
                            except UnicodeDecodeError:
                                uploaded_file.seek(0)
                                reader = pd.read_csv(uploaded_file, chunksize=chunk_size, encoding='ISO-8859-1', on_bad_lines='skip')
                                
                            for chunk in reader:
                                if chunk.empty: continue
                                if len(chunk.columns) > MAX_COLUMNS:
                                    chunk = chunk.iloc[:, :MAX_COLUMNS]
                                if first_chunk:
                                    chunk.to_sql('user_data', conn, if_exists='replace', index=False)
                                    st.session_state.dataset_columns = ", ".join(chunk.columns.tolist())
                                    if len(chunk.columns) == MAX_COLUMNS:
                                        st.warning(f"Trimmed to {MAX_COLUMNS} columns.")
                                    first_chunk = False
                                else:
                                    chunk.to_sql('user_data', conn, if_exists='append', index=False)
                            progress_text.empty()
                        except pd.errors.EmptyDataError:
                            st.error("The uploaded CSV file is empty.")
                        
                    elif data_type == "JSON":
                        try:
                            df = pd.read_json(uploaded_file)
                            if df.empty:
                                st.error("The uploaded JSON file is empty.")
                            else:
                                if len(df.columns) > MAX_COLUMNS:
                                    df = df.iloc[:, :MAX_COLUMNS]
                                df.to_sql('user_data', conn, if_exists='replace', index=False)
                                st.session_state.dataset_columns = ", ".join(df.columns.tolist())
                        except ValueError as e:
                            st.error(f"Invalid JSON format: {e}")
                        
                    elif data_type == "PDF (Document)":
                        try:
                            import pdfplumber
                            pdf_rows = []
                            with pdfplumber.open(uploaded_file) as pdf:
                                for page_num, page in enumerate(pdf.pages, 1):
                                    text = page.extract_text() or ""
                                    # Chunk at line boundaries to avoid breaking mid-sentence
                                    chunks = []
                                    if text:
                                        current = ""
                                        for line in text.splitlines(keepends=True):
                                            if len(current) + len(line) > 1000 and current:
                                                chunks.append(current.strip())
                                                current = line
                                            else:
                                                current += line
                                        if current.strip():
                                            chunks.append(current.strip())
                                    for chunk_idx, chunk_text in enumerate(chunks):
                                        pdf_rows.append({
                                            "page_number": page_num,
                                            "chunk_index": chunk_idx,
                                            "text_content": chunk_text.strip()
                                        })
                            
                            if pdf_rows:
                                df = pd.DataFrame(pdf_rows)
                                df.to_sql('user_data', conn, if_exists='replace', index=False)
                                st.session_state.dataset_columns = "page_number, chunk_index, text_content"
                            else:
                                st.warning("No text could be extracted from this PDF. It might be scanned or encrypted.")
                        except Exception as e:
                            st.error(f"Failed to read PDF file: {e}")
     
                    elif data_type == "PPTX (Slides)":
                        try:
                            from pptx import Presentation
                            pptx_rows = []
                            prs = Presentation(uploaded_file)
                            for slide_num, slide in enumerate(prs.slides, 1):
                                slide_text = []
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        slide_text.append(shape.text)
                                
                                full_slide_text = "\n".join(slide_text).strip()
                                if full_slide_text:
                                    # For slides, we typically keep the whole slide content together unless it's massive
                                    pptx_rows.append({
                                        "page_number": slide_num,
                                        "chunk_index": 0,
                                        "text_content": full_slide_text
                                    })
                            
                            if pptx_rows:
                                df = pd.DataFrame(pptx_rows)
                                df.to_sql('user_data', conn, if_exists='replace', index=False)
                                st.session_state.dataset_columns = "page_number, chunk_index, text_content"
                            else:
                                st.warning("No text could be found in these slides.")
                        except Exception as e:
                            st.error(f"Failed to read PPTX file: {e}")

                    elif data_type == "DOCX (Document)":
                        try:
                            import docx
                            doc = docx.Document(uploaded_file)
                            docx_rows = []
                            full_text = []
                            for para in doc.paragraphs:
                                if para.text.strip():
                                    full_text.append(para.text.strip())
                            
                            joined_text = "\n".join(full_text)
                            # Chunk at paragraph boundaries to avoid breaking mid-sentence
                            chunks = []
                            current = ""
                            for line in joined_text.splitlines(keepends=True):
                                if len(current) + len(line) > 1000 and current:
                                    chunks.append(current.strip())
                                    current = line
                                else:
                                    current += line
                            if current.strip():
                                chunks.append(current.strip())
                            for chunk_idx, chunk_text in enumerate(chunks):
                                docx_rows.append({
                                    "page_number": 1, # Word doesn't have reliable page numbers via this method
                                    "chunk_index": chunk_idx,
                                    "text_content": chunk_text.strip()
                                })
                            
                            if docx_rows:
                                df = pd.DataFrame(docx_rows)
                                df.to_sql('user_data', conn, if_exists='replace', index=False)
                                st.session_state.dataset_columns = "page_number, chunk_index, text_content"
                            else:
                                st.warning("This document appears to be empty.")
                        except Exception as e:
                            st.error(f"Failed to read Word document: {e}")

                except Exception as e:
                    st.error(f"Processing failed: {e}")
                finally:
                    if conn:
                        conn.close()
                    # Mark as ingested so we don't repeat this on next rerun
                    st.session_state.last_ingested_file_id = file_unique_id
                    
                    # Trigger proactive EDA for structured data (CSV/JSON only)
                    if data_type in ("CSV (Spreadsheet)", "JSON"):
                        st.session_state.run_proactive_eda = True
                    
                    # Switch mode if not already there
                    if st.session_state.confirmed_mode != "Local (Uploaded CSV - SQLite)":
                        st.session_state.active_mode = "Local (Uploaded CSV - SQLite)"
                        st.session_state.confirmed_mode = "Local (Uploaded CSV - SQLite)"
                        st.rerun()

            # Always show success and preview if table exists (even if already ingested)
            check_conn = None
            try:
                check_conn = sqlite3.connect('hackathon_database.db')
                check_table = pd.read_sql_query("SELECT count(*) FROM sqlite_master WHERE type='table' AND name='user_data'", check_conn)
                if check_table.iloc[0, 0] > 0:
                    st.success("Database Active!")
                    with st.expander("Preview Local Data"):
                        preview_df = pd.read_sql_query("SELECT * FROM user_data LIMIT 5", check_conn)
                        st.dataframe(preview_df)
            except Exception as e:
                st.error(f"Could not load database preview: {e}")
            finally:
                if check_conn:
                    check_conn.close()
                
        st.divider()
        
        # VISION PROCESSOR — Images + Videos
        st.markdown("### Vision")
        vision_file = st.file_uploader(
            "Upload Image or Video for Analysis", 
            type=["png", "jpg", "jpeg", "mp4"], 
            key="vision_uploader"
        )
        
        if vision_file is not None:
            file_type = vision_file.type or ""

            if file_type in ["image/png", "image/jpeg", "image/jpg"]:
                # IMAGE handling
                img_bytes = vision_file.read()
                MAX_IMAGE_BYTES = 10 * 1024 * 1024  # 10 MB
                if len(img_bytes) > MAX_IMAGE_BYTES:
                    st.error("Image exceeds 10 MB limit. Please upload a smaller image.")
                    img_bytes = None
                if img_bytes:
                    st.session_state.uploaded_image_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    st.session_state.uploaded_video_bytes = None
                    st.image(img_bytes, caption="Uploaded Image", use_container_width=True)

                    if st.button("🔍 Analyze Image", use_container_width=True):
                        with st.spinner("Sending to Groq Vision model..."):
                            result = agent.execute_vision_tool(
                                st.session_state.uploaded_image_b64,
                                "Describe this image in detail. Extract any text, data, charts, or key visual elements."
                            )
                            st.markdown("**Analysis Result:**")
                            st.markdown(result)
                            st.session_state.messages.append({"role": "assistant", "content": f"🖼️ **Vision Analysis:**\n\n{result}"})
                        
            elif file_type == "video/mp4":
                # VIDEO handling
                video_bytes = vision_file.read()
                st.session_state.uploaded_video_bytes = video_bytes
                st.session_state.uploaded_image_b64 = None
                st.video(video_bytes)
                
                if st.button("🎬 Analyze Video", use_container_width=True):
                    with st.spinner("Extracting frames & analyzing with Vision model..."):
                        result = agent.execute_video_tool(
                            video_bytes,
                            "Analyze this video frame. Describe what you see including objects, actions, text, and scenery."
                        )
                        st.markdown("**Video Analysis Result:**")
                        st.markdown(result)
                        st.session_state.messages.append({"role": "assistant", "content": f"🎬 **Video Analysis:**\n\n{result}"})
        else:
            st.session_state.uploaded_image_b64 = None
            st.session_state.uploaded_video_bytes = None
        
        st.caption("Powered by Groq Vision (Llama 3.2 90B) • Supports PNG, JPG, MP4")

    # TAB 3: INTEGRATIONS
    with tab_tools:
        st.markdown("### Integrations")
        st.session_state.contact_target = st.text_input(
            "Target Phone / Email",
            value=st.session_state.contact_target,
            placeholder="+1 (555) 000-0000"
        )
        st.caption("Used by the Agent to send automated reports. Note: message sending is simulated — no actual messages are delivered.")

# --- 4. THE ANALYST DASHBOARD ---
_dash_conn = None
try:
    _dash_conn = sqlite3.connect('hackathon_database.db')
    check_df = pd.read_sql_query("SELECT 1 FROM user_data LIMIT 1", _dash_conn)

    if not check_df.empty:
        with st.expander("📊 Open Analyst Dashboard"):
            columns_df = pd.read_sql_query("PRAGMA table_info(user_data)", _dash_conn)
            all_columns = columns_df['name'].tolist()
            sample_df = pd.read_sql_query("SELECT * FROM user_data LIMIT 100", _dash_conn)
            numeric_columns = sample_df.select_dtypes(include=['number']).columns.tolist()

            if not all_columns or not numeric_columns:
                st.info("Insufficient structured numerical data for visualization.")
            else:
                col1, col2, col3 = st.columns(3)
                with col1: chart_type = st.selectbox("Chart", ["Bar Chart", "Line Chart", "Scatter Plot"])
                with col2: x_axis = st.selectbox("X-Axis", all_columns)
                with col3: y_axis = st.selectbox("Y-Axis", numeric_columns)

                if st.button("Generate Visualization", use_container_width=True):
                    try:
                        # Use a fresh connection for the chart query to avoid state issues
                        with sqlite3.connect('hackathon_database.db') as chart_conn:
                            query = f"SELECT `{x_axis}`, `{y_axis}` FROM user_data LIMIT 1000"
                            chart_data = pd.read_sql_query(query, chart_conn).set_index(x_axis)

                            if chart_data.empty:
                                st.warning("No data found for the selected axes.")
                            else:
                                if chart_type == "Bar Chart": st.bar_chart(chart_data)
                                elif chart_type == "Line Chart": st.line_chart(chart_data)
                                elif chart_type == "Scatter Plot":
                                    scatter_df = pd.read_sql_query(query, chart_conn)
                                    st.scatter_chart(scatter_df, x=x_axis, y=y_axis)
                    except Exception as e:
                        st.error(f"Visualization error: Failed to map data for axes properly. {e}")
except Exception:
    pass  # No database yet — dashboard is intentionally hidden
finally:
    if _dash_conn:
        _dash_conn.close()

st.write("") # Clean spacing

# --- PROACTIVE EDA HELPER ---
def compute_eda_summary():
    """Computes basic EDA stats from the user_data table for the proactive briefing."""
    eda_conn = None
    try:
        eda_conn = sqlite3.connect('hackathon_database.db', timeout=10)
        
        # Row count
        row_count = pd.read_sql_query("SELECT COUNT(*) as cnt FROM user_data", eda_conn).iloc[0, 0]
        
        # Column info
        col_info = pd.read_sql_query("PRAGMA table_info(user_data)", eda_conn)
        column_names = col_info['name'].tolist()
        column_types = col_info['type'].tolist()
        
        # Sample for numeric analysis
        sample_df = pd.read_sql_query("SELECT * FROM user_data LIMIT 100", eda_conn)
        numeric_cols = sample_df.select_dtypes(include=['number']).columns.tolist()
        
        # Build summary
        lines = [
            f"**Rows:** {row_count}",
            f"**Columns ({len(column_names)}):** {', '.join(column_names)}",
            f"**Column types:** {', '.join(f'{n} ({t})' for n, t in zip(column_names, column_types))}",
        ]
        
        if numeric_cols:
            lines.append(f"**Numeric columns:** {', '.join(numeric_cols)}")
            for col in numeric_cols[:5]:  # Limit to first 5 numeric columns
                stats = sample_df[col].describe()
                missing = sample_df[col].isnull().sum()
                lines.append(
                    f"  - {col}: min={stats.get('min', 'N/A')}, max={stats.get('max', 'N/A')}, "
                    f"mean={stats.get('mean', 'N/A'):.2f}, missing={missing}"
                )
        else:
            lines.append("**Numeric columns:** None detected (text-based dataset)")
        
        # Missing values overview
        total_missing = sample_df.isnull().sum()
        cols_with_missing = [(c, v) for c, v in total_missing.items() if v > 0]
        if cols_with_missing:
            lines.append(f"**Columns with missing values (in first 100 rows):** {', '.join(f'{c} ({v})' for c, v in cols_with_missing)}")
        else:
            lines.append("**Missing values:** None detected in sample")
        
        return "\n".join(lines)
    except Exception as e:
        return f"Could not compute EDA: {e}"
    finally:
        if eda_conn:
            eda_conn.close()



# --- 5. THE CHAT INTERFACE ---

# CSS for thinking display — theme-aware
if is_cloud:
    think_border = "#8fa8c0"
    think_bg = "rgba(143,168,192,0.08)"
    ring_color = "#8fa8c0"
else:
    think_border = "#c8a87a"
    think_bg = "rgba(200,168,122,0.09)"
    ring_color = "#c8a87a"

st.markdown(f"""
<style>
/* ── Thinking box ── */
.thinking-box {{
    background: {think_bg};
    border-left: 2px solid {think_border};
    border-radius: 0 10px 10px 0;
    padding: 14px 18px;
    margin: 10px 0;
    font-size: 0.88em;
    line-height: 1.6;
    font-family: 'Zen Kaku Gothic New', sans-serif;
}}
.thinking-header {{
    font-family: 'Shippori Mincho', serif;
    font-weight: 600;
    font-size: 0.95em;
    margin-bottom: 6px;
    display: flex;
    align-items: center;
    gap: 10px;
    color: {think_border};
    letter-spacing: 0.04em;
}}

/* ── Think-ring: meditative rotating arc ── */
.think-ring {{
    display: inline-block;
    width: 18px;
    height: 18px;
    border-radius: 50%;
    border: 2px solid transparent;
    border-top-color: {ring_color};
    border-right-color: {ring_color}3a;
    animation: think-rotate 3s linear infinite;
    vertical-align: middle;
}}
@keyframes think-rotate {{
    0%   {{ transform: rotate(0deg); }}
    100% {{ transform: rotate(360deg); }}
}}

/* ── Tool badge ── */
.tool-badge {{
    display: inline-block;
    padding: 3px 12px;
    border-radius: 20px;
    font-size: 0.78em;
    font-family: 'Zen Kaku Gothic New', sans-serif;
    font-weight: 500;
    letter-spacing: 0.08em;
}}
</style>
""", unsafe_allow_html=True)

def render_chat_message(msg):
    # Avatar SVGs
    if is_cloud:
        # Snow Leopard Avatars
        user_avatar = f"data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ccircle cx='50' cy='50' r='50' fill='%23dde8f2'/%3E%3Ctext x='50' y='53' font-family='Shippori Mincho, serif' font-size='44' font-weight='500' fill='%233a5068' text-anchor='middle' dominant-baseline='middle'%3EU%3C/text%3E%3C/svg%3E"
        agent_avatar = f"data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ccircle cx='50' cy='50' r='50' fill='%238fa8c0'/%3E%3Ccircle cx='35' cy='35' r='6' fill='white'/%3E%3Ccircle cx='65' cy='40' r='5' fill='white'/%3E%3Ccircle cx='40' cy='65' r='7' fill='white'/%3E%3Ccircle cx='60' cy='60' r='4' fill='white'/%3E%3C/svg%3E"
    else:
        # Local Avatars
        user_avatar = f"data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ccircle cx='50' cy='50' r='50' fill='%23ede0cc'/%3E%3Ctext x='50' y='53' font-family='Shippori Mincho, serif' font-size='44' font-weight='500' fill='%235c3d1e' text-anchor='middle' dominant-baseline='middle'%3EU%3C/text%3E%3C/svg%3E"
        agent_avatar = f"data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ccircle cx='50' cy='50' r='50' fill='%23c8a87a'/%3E%3Cpath d='M 30,50 A 20,20 0 1,1 60,65' fill='none' stroke='white' stroke-width='4' stroke-linecap='round'/%3E%3C/svg%3E"

    avatar_to_use = user_avatar if msg["role"] == "user" else agent_avatar

    with st.chat_message(msg["role"], avatar=avatar_to_use):
        st.markdown(msg["content"])
        
        # Display references if any
        if 'references' in msg and msg['references']:
            with st.expander("Sources & Verification", expanded=False):
                for ref in msg['references']:
                    st.markdown(f"- {ref}")
        
        # Display thinking block if it exists
        if 'thinking' in msg and msg['thinking']:
            with st.expander("Agent Thought Process", expanded=False):
                st.markdown(f"```python\n{msg['thinking']}\n```")

        # Display chart if generated
        if 'chart' in msg and msg['chart'] is not None:
            chart_data = msg['chart']
            x_col = chart_data.columns[0]
            st.bar_chart(chart_data.set_index(x_col))

for msg in st.session_state.messages:
    render_chat_message(msg)

# --- PROACTIVE EDA RENDERING ---
if st.session_state.run_proactive_eda:
    st.session_state.run_proactive_eda = False  # Clear flag immediately
    
    eda_summary = compute_eda_summary()
    
    # Determine avatar for the briefing message
    if is_cloud:
        briefing_avatar = f"data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ccircle cx='50' cy='50' r='50' fill='%238fa8c0'/%3E%3Ccircle cx='35' cy='35' r='6' fill='white'/%3E%3Ccircle cx='65' cy='40' r='5' fill='white'/%3E%3Ccircle cx='40' cy='65' r='7' fill='white'/%3E%3Ccircle cx='60' cy='60' r='4' fill='white'/%3E%3C/svg%3E"
    else:
        briefing_avatar = f"data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ccircle cx='50' cy='50' r='50' fill='%23c8a87a'/%3E%3Cpath d='M 30,50 A 20,20 0 1,1 60,65' fill='none' stroke='white' stroke-width='4' stroke-linecap='round'/%3E%3C/svg%3E"
    
    with st.chat_message("assistant", avatar=briefing_avatar):
        briefing_answer = st.write_stream(
            agent.stream_proactive_briefing(eda_summary)
        )
    
    st.session_state.messages.append({"role": "assistant", "content": briefing_answer})


# Use voice transcript if available, otherwise use chat input
chat_prompt = None
if st.session_state.voice_transcript:
    chat_prompt = st.session_state.voice_transcript
    st.session_state.voice_transcript = None  # Clear after use

if manual_prompt := st.chat_input("Ask me to analyze data, identify an image, or send a report..."):
    chat_prompt = manual_prompt

if chat_prompt:
    st.session_state.messages.append({"role": "user", "content": chat_prompt})
    render_chat_message({"role": "user", "content": chat_prompt})

    # Avatar SVG logic for live agent response
    if is_cloud:
        current_agent_avatar = f"data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ccircle cx='50' cy='50' r='50' fill='%238fa8c0'/%3E%3Ccircle cx='35' cy='35' r='6' fill='white'/%3E%3Ccircle cx='65' cy='40' r='5' fill='white'/%3E%3Ccircle cx='40' cy='65' r='7' fill='white'/%3E%3Ccircle cx='60' cy='60' r='4' fill='white'/%3E%3C/svg%3E"
    else:
        current_agent_avatar = f"data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ccircle cx='50' cy='50' r='50' fill='%23c8a87a'/%3E%3Cpath d='M 30,50 A 20,20 0 1,1 60,65' fill='none' stroke='white' stroke-width='4' stroke-linecap='round'/%3E%3C/svg%3E"

    with st.chat_message("assistant", avatar=current_agent_avatar):
        # Thinking UI placeholder
        thinking_placeholder = st.empty()
        
        # Enso Loader SVG
        accent_color = ACCENT_U if is_cloud else ACCENT_U_LOCAL
        enso_svg = f"""<svg width='18' height='18' viewBox='0 0 40 40' style='display:inline-block; vertical-align:-3px; margin-right:8px;'>
                        <path class='enso-draw' d='M 20,4 A 16,16 0 1,1 5,28' fill='none' stroke='{accent_color}' stroke-width='3' stroke-linecap='round' pathLength='100'/>
                    </svg>"""

    thinking_html = f"""
    <div class="thinking-box">
        <div class="thinking-header">
            {enso_svg} Agent Thought
        </div>
        <div class="thinking-content">
            <span style="opacity:0.7">Initiating thought sequence...</span>
        </div>
    </div>
    """
    thinking_placeholder.markdown(thinking_html, unsafe_allow_html=True)
        
    raw_decision, tool_name, tool_result, needs_synthesis = agent.ask_agent(
        chat_prompt, 
        st.session_state.messages, 
        st.session_state.dataset_columns,
        st.session_state.confirmed_mode,
        temperature=st.session_state.temperature,
        contact_target=st.session_state.contact_target if st.session_state.contact_target else None,
        has_image=st.session_state.uploaded_image_b64 is not None,
        has_video=st.session_state.uploaded_video_bytes is not None,
        video_bytes=st.session_state.uploaded_video_bytes
    )
    
    # Parse and display thinking
    agent_thought = "Processing..."
    tool_used = tool_name or "Unknown"
    try:
        start_idx = raw_decision.find('{')
        end_idx = raw_decision.rfind('}')
        if start_idx != -1 and end_idx > start_idx:
            clean_json = raw_decision[start_idx:end_idx+1]
            decision_dict = json.loads(clean_json)
            agent_thought = decision_dict.get("thought", "Processing...")
            tool_used = decision_dict.get("tool", tool_name or "Unknown")
    except (json.JSONDecodeError, ValueError):
        pass  # Keep defaults; raw_decision was already used by ask_agent

    # Escape LLM-generated text before embedding in HTML to prevent XSS
    agent_thought_safe = html.escape(agent_thought)

    # Tool badge formatting and human-readable names
    tool_info = {
        "query_live_data": ("☁️", "Searching Database", "#2196f3", "rgba(33,150,243,0.12)"),
        "query_local_data": ("💾", "Analyzing Files", "#4caf50", "rgba(76,175,80,0.12)"),
        "analyze_image": ("🖼️", "Examining Image", "#9c27b0", "rgba(156,39,176,0.12)"),
        "analyze_video": ("🎬", "Reviewing Video", "#e91e63", "rgba(233,30,99,0.12)"),
        "send_message": ("📨", "Sending Message", "#ff9800", "rgba(255,152,0,0.12)"),
        "direct_response": ("💬", "Generating Answer", "#607d8b", "rgba(96,125,139,0.12)"),
    }
    icon, label, badge_color, badge_bg = tool_info.get(tool_used, ("⚙️", "Processing", "#607d8b", "rgba(96,125,139,0.12)"))

    thinking_placeholder.markdown(
        f'<div class="thinking-box">'
        f'<div class="thinking-header">Agent Thought</div>'
        f'<p style="margin:6px 0 12px 0;opacity:0.88; font-size: 0.95rem;">{agent_thought_safe}</p>'
        f'<span class="tool-badge" style="background:{badge_bg};color:{badge_color};border:0.5px solid {badge_color}40;">'
        f'{icon}&nbsp;{label}</span>'
        f'</div>',
        unsafe_allow_html=True
    )
    
    # --- PHASE 2: RESPONSE (streamed) ---
    import time
    
    if needs_synthesis:
        # Stream the synthesis response token by token
        final_answer = st.write_stream(
            agent.stream_synthesis(chat_prompt, tool_result, st.session_state.temperature)
        )
    else:
        # For non-database tools, simulate typing for the result
        def type_result(text):
            """Generator that yields small chunks to create typewriter effect."""
            words = text.split(' ')
            for i, word in enumerate(words):
                yield word + (' ' if i < len(words) - 1 else '')
                time.sleep(0.02)
        
        final_answer = st.write_stream(type_result(tool_result))
            
    st.session_state.messages.append({"role": "assistant", "content": final_answer})

# --- 6. EXPORT REPORT ENGINE ---
if st.session_state.messages:
    st.divider()
    
    # Build the report from recent messages
    def build_report():
        report_lines = [
            "# 📊 Acuman AI — Analysis Report",
            f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            f"**Mode:** {st.session_state.confirmed_mode}",
            f"**Creativity:** {st.session_state.temperature}",
            "",
            "---",
            ""
        ]
        for msg in st.session_state.messages:
            role = "🧑 User" if msg["role"] == "user" else "🤖 Acuman"
            report_lines.append(f"### {role}")
            report_lines.append(msg["content"])
            report_lines.append("")
        
        report_lines.append("---")
        report_lines.append("*Report generated by Acuman: Hybrid Neural Agent*")
        return "\n".join(report_lines)
    
    report_md = build_report()
    st.download_button(
        label="📥 Export Report",
        data=report_md,
        file_name=f"acuman_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md",
        mime="text/markdown",
        use_container_width=True
    )